#!/usr/bin/env python3
"""Build the boxBot / Smart Home 2.0 deck for Alarm.com CX PMs.

Minimal slides to accompany a spoken talking track. Visual language matches
boxBot's own display theme: warm amber/coral on a near-black warm background.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")

# ---- palette (sampled from boxBot's 'boxbot' display theme) -----------------
BG       = RGBColor(0x17, 0x14, 0x11)   # warm near-black
PANEL    = RGBColor(0x26, 0x20, 0x1B)   # card surface
HAIRLINE = RGBColor(0x3C, 0x34, 0x2E)
AMBER    = RGBColor(0xE8, 0x91, 0x5B)   # primary accent
AMBER_DK = RGBColor(0xD2, 0x7A, 0x48)
INK      = RGBColor(0xF1, 0xEC, 0xE5)   # primary text
MUTED    = RGBColor(0x9A, 0x90, 0x86)   # secondary text
FAINT    = RGBColor(0x6B, 0x62, 0x59)

FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- helpers ----------------------------------------------------------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element; sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=None, line_spacing=None):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, tracking)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for (txt, size, color, bold, *rest) in para:
            r = p.add_run(); r.text = txt
            f = r.font
            f.name = FONT; f.size = Pt(size); f.bold = bold
            f.color.rgb = color
            if rest and rest[0]:
                _track(r, rest[0])
    return tb


def _track(run, pts):
    """Letter-spacing in points (x100 in EMU-ish 'spc' attr)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(pts * 100)))


def bar(s, x, y, w, h, color=AMBER):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def panel(s, x, y, w, h, fill=PANEL, line=None):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.adjustments[0] = 0.045
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line:
        r.line.color.rgb = line; r.line.width = Pt(1)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def kicker(s, x, y, label, color=AMBER):
    """Small uppercase section label with a short rule above it."""
    bar(s, x, y, Inches(0.5), Pt(3), color)
    text(s, x, y + Pt(10), Inches(8), Inches(0.4),
         [[(label.upper(), 13, color, True, 3)]])


def placeholder(s, x, y, w, h, label, sub=""):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.adjustments[0] = 0.04
    p.fill.solid(); p.fill.fore_color.rgb = PANEL
    p.line.color.rgb = AMBER_DK; p.line.width = Pt(1.25)
    p.line.dash_style = None
    p.shadow.inherit = False
    # dashed border
    ln = p.line._get_or_add_ln()
    d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(d)
    runs = [[("▶  " + label, 18, AMBER, True)]]
    if sub:
        runs.append([(sub, 12.5, MUTED, False)])
    text(s, x, y, w, h, runs, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, space_after=6)
    return p


def image_contain(s, path, x, y, w, h, border=True):
    iw, ih = Image.open(path).size
    ar_box = w / h
    ar_img = iw / ih
    if ar_img > ar_box:
        nw = w; nh = int(w / ar_img)
    else:
        nh = h; nw = int(h * ar_img)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    if border:
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx - Emu(9144),
                               ny - Emu(9144), nw + Emu(18288), nh + Emu(18288))
        b.adjustments[0] = 0.03
        b.fill.background(); b.line.color.rgb = HAIRLINE; b.line.width = Pt(1)
        b.shadow.inherit = False
    s.shapes.add_picture(path, nx, ny, nw, nh)


def notes(s, txt, tech=None):
    body = txt.strip()
    if tech:
        body += ("\n\n" + "—" * 28 + "\nTECHNICAL IMPLEMENTATION  (optional depth — weave in as you like)\n"
                 + "—" * 28 + "\n" + tech.strip())
    s.notes_slide.notes_text_frame.text = body


MX = Inches(0.92)   # left margin


# =============================================================================
# 1 — TITLE
# =============================================================================
s = slide()
bar(s, MX, Inches(2.18), Inches(0.7), Pt(4), AMBER)
text(s, MX, Inches(2.45), Inches(9.5), Inches(0.5),
     [[("boxBot", 17, AMBER, True, 4)]])
text(s, MX, Inches(2.95), Inches(11), Inches(2.2),
     [[("Smart Home 2.0", 58, INK, True)],
      [("An AI-native approach to the connected home", 24, MUTED, False)]],
     space_after=14)
text(s, MX, Inches(6.35), Inches(11), Inches(0.6),
     [[("A working prototype — and a different way to think about the platform", 15, FAINT, False)]])
placeholder(s, Inches(8.7), Inches(0.7), Inches(3.7), Inches(6.1),
            "BB PHOTO", "the device, on the counter")
notes(s, """
Frame the room: we're not here to look at a new feature. We're here to look at a
different shape for the whole product. The prototype on screen runs an agent
harness I built on weekends — it's in daily use in my house. Everything I show is
real. The goal today is to open up a different way of thinking about what the
smart home platform could be in the AI age.
""", tech="""
• Hardware: Raspberry Pi 5 (8 GB) + AI HAT+ (Hailo-8L NPU, 13 TOPS) in a wooden box.
  Camera (IMX708 NoIR), ReSpeaker 4-mic USB array, 7" 1024x600 LCD, speaker.
• One long-lived Python process (src/boxbot/core/main.py). The conversation loop
  runs on the Claude Agent SDK backend (agent.backend="claude_agent_sdk";
  core/agent.py::_agent_loop_sdk, bridged by core/agent_sdk_adapter.py — boxBot
  tools are wrapped into an in-process MCP server).
• Two models (core/config.py): BOXBOT_MODEL_LARGE = conversation/reasoning;
  BOXBOT_MODEL_SMALL = Haiku for high-frequency, low-stakes work (web-search
  content firewall, memory rerank, photo tagging).
• Everything is Pi-local. The only data leaving the box is explicit API calls
  (Claude, ElevenLabs STT/TTS, WhatsApp, weather). No telemetry; the only open
  inbound port is the WhatsApp webhook.
""")

# =============================================================================
# 2 — THE THESIS (one line)
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "The thesis")
text(s, MX, Inches(2.5), Inches(11.5), Inches(3.2),
     [[("The platform of the next decade is ", 38, INK, False),
       ("not", 38, MUTED, False),
       (" our app with an assistant added on.", 38, INK, False)],
      [("It ", 38, INK, False),
       ("is", 38, AMBER, True),
       (" the assistant, with the home built around it.", 38, INK, True)]],
     space_after=18, line_spacing=1.05)
notes(s, """
This is the whole talk in one sentence. Today, almost everyone in the industry —
us included — is bolting an assistant onto an existing app. A voice on top of the
same menus. The bet here is the opposite: the assistant is the platform, and the
app, the panel, the settings all become things the agent operates on the user's
behalf. Hold that inversion in your head for the next ten minutes.
""", tech="""
• The inversion is literal in the codebase: there is no settings UI, no rules
  engine, no React/Swift app. "App behavior" is just text files the agent reads
  and edits:
    – skills  → skills/<name>/SKILL.md  (instructions)
    – displays → declarative JSON specs  (layout)
    – standing behavior → data/memory/system.md  (always-loaded household facts)
• The entire "frontend" the agent is given = 10 always-loaded tools + the `bb`
  sandbox SDK (src/boxbot/sdk/). That's the surface area we maintain.
• So "adding a feature" = adding/editing a text file, not shipping a build. That
  is the structural claim the rest of the talk makes concrete.
""")

# =============================================================================
# 3 — TODAY'S SMART HOME (the friction)
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "Where we are today")
text(s, MX, Inches(1.7), Inches(11.5), Inches(1.0),
     [[("To get what you want, you go to the app.", 33, INK, True)]])
items = ["Fixed widgets", "Notification rules", "Settings pages",
         "Automation wizards", "Device menus", "Help articles"]
cols = 3
cw = Inches(3.7); gap = Inches(0.28)
cardh = Inches(1.15)
x0 = MX; y0 = Inches(3.05)
for i, it in enumerate(items):
    r = i // cols; c = i % cols
    x = x0 + c * (cw + gap)
    y = y0 + r * (cardh + gap)
    panel(s, x, y, cw, cardh, fill=PANEL)
    text(s, x + Inches(0.3), y, cw - Inches(0.5), cardh,
         [[(it, 19, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, MX, Inches(6.75), Inches(11.5), Inches(0.5),
     [[("Every capability is a surface the user has to find, learn, and configure.", 16, MUTED, False)]])
notes(s, """
Walk the grid. This is the modern smart-home app — and it's a good app. But notice
the pattern: every single thing the product can do is a surface. A screen you have
to discover, a menu you have to learn, a rule you have to assemble. The intelligence
lives in the user's head — they have to translate what they want into the app's
vocabulary. The richer the product gets, the more surfaces there are to wrangle.
Keep this grid in mind; we're going to make it disappear.
""", tech="""
• Contrast to draw if useful: in a conventional stack each tile on this grid is
  code we build and own — a settings screen, a rule row in a DB, a notification
  dispatcher service, a feature flag, a help-doc CMS, plus the test matrix for all
  of it.
• In the agent model those collapse into four reusable primitives, none of which
  the user navigates:
    – skills (instructions)           – integrations (data pipes, manifest+script)
    – displays (declarative specs)    – triggers (wake conditions, core/scheduler.py)
• There's nothing to navigate because there's no UI to navigate. The agent IS the
  interface; the "surfaces" become things it operates, not things the user learns.
""")

# =============================================================================
# 4 — THE SAME JOB, TWO WAYS (setup)
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "A concrete example")
text(s, MX, Inches(2.5), Inches(11.5), Inches(3.0),
     [[("“Text my wife whenever I add", 40, INK, True)],
      [("something to the family calendar.”", 40, INK, True)]],
     space_after=8, line_spacing=1.06)
text(s, MX, Inches(5.0), Inches(11), Inches(0.6),
     [[("A real request from my house. Let's do it both ways.", 19, AMBER, False)]])
notes(s, """
Here's a real job-to-be-done from my own family. My wife was the household
scheduler; she wanted to stay in the loop when I add things. Simple human request.
Let's watch what it takes to satisfy it in today's model versus the agent model.
This is the heart of the talk — the next three slides are the payoff.
""", tech="""
• Worth knowing for Q&A: the calendar is an INTEGRATION, not a skill —
  integrations/calendar/{manifest.yaml, script.py}. The manifest declares inputs
  (action: list_upcoming_events / create_event / update_event / delete_event),
  outputs, a required secret (GOOGLE_CALENDAR_TOKEN_JSON), and a 20s timeout.
• script.py runs sandboxed; it handles Google OAuth token refresh and persists the
  rotated token back via bb.secrets. The agent reaches it with:
    bb.integrations.get("calendar", action="list_upcoming_events", max_results=5)
• So "the calendar" is already a clean, callable data pipe before any of the
  notification logic exists. That separation is what makes the next slide a
  one-liner.
""")

# =============================================================================
# 5 — THE OLD WAY (rule engine)
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.8), "The old way", AMBER_DK)
text(s, MX, Inches(1.45), Inches(11.5), Inches(0.7),
     [[("Open settings. Build the rule.", 30, INK, True)]])
# rule-builder mock
rx, ry, rw = MX, Inches(2.55), Inches(11.5)
panel(s, rx, ry, rw, Inches(3.2), fill=PANEL, line=HAIRLINE)
rows = [
    ("SEND", "SMS  /  Push  /  Email   ?"),
    ("TO", "User 1, User 2, …   (pick recipients)"),
    ("WHEN", "Service = Calendar   ·   Condition = item added"),
    ("DURING", "Timeframe / quiet hours / days of week   ?"),
]
iy = ry + Inches(0.32)
for label, val in rows:
    text(s, rx + Inches(0.45), iy, Inches(2.2), Inches(0.5),
         [[(label, 14, AMBER, True, 2)]])
    text(s, rx + Inches(2.9), iy, rw - Inches(3.2), Inches(0.5),
         [[(val, 17, MUTED, False)]])
    iy += Inches(0.72)
text(s, MX, Inches(6.25), Inches(11.5), Inches(0.8),
     [[("Pick the channel. Pick the people. Find the trigger. Set the conditions. ", 16, MUTED, False),
       ("And someone had to build that rule engine first.", 16, INK, True)]])
notes(s, """
The old way. Even when the app supports it, the user has to become a rules
programmer for a moment: choose a channel, choose recipients, find the right
trigger, scope the timeframe. And step back — for ANY of this to exist, our
engineering org had to design and build and maintain that whole rule-builder:
the UI, the condition logic, the delivery backend, the test matrix. Multiply that
by every feature. That's the cost structure of the current model.
""", tech="""
• The point to land: this column is the part we DON'T build in the agent model.
  In today's stack, shipping this one feature means designing and maintaining:
    – a rule schema + storage          – a builder UI (channel/recipient/condition pickers)
    – a condition evaluator            – a multi-channel delivery service (SMS/push/email)
    – quiet-hours / timeframe logic    – the regression test matrix for all of it
• Every one of those is durable engineering surface that has to be kept alive
  across OS versions, channel API changes, and edge cases — forever, per feature.
• The next slide deletes this entire column.
""")

# =============================================================================
# 6 — THE NEW WAY (one text -> one line)
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.8), "The new way")
text(s, MX, Inches(1.45), Inches(11.5), Inches(0.7),
     [[("She just tells boxBot.", 30, INK, True)]])
# chat bubble
bx, by, bw = MX, Inches(2.5), Inches(7.2)
panel(s, bx, by, bw, Inches(1.15), fill=PANEL)
text(s, bx + Inches(0.4), by, bw - Inches(0.7), Inches(1.15),
     [[("“Hey BB, text me whenever Jacob adds something to the calendar.”", 18, INK, False)]],
     anchor=MSO_ANCHOR.MIDDLE)
# arrow
text(s, MX, Inches(3.95), Inches(11), Inches(0.5),
     [[("↓   boxBot edits its calendar skill", 16, AMBER, True)]])
# code-ish line
panel(s, MX, Inches(4.55), Inches(11.5), Inches(1.15), fill=RGBColor(0x10,0x0D,0x0A), line=HAIRLINE)
text(s, MX + Inches(0.45), Inches(4.55), Inches(11), Inches(1.15),
     [[("calendar/SKILL.md  +", 13, FAINT, False, 1)],
      [("On calendar add → text Sarah a one-line summary.", 19, AMBER, True)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=6)
text(s, MX, Inches(6.15), Inches(11.5), Inches(0.7),
     [[("One sentence in. One line added. ", 18, INK, True),
       ("Now every calendar add texts her — no settings, no rule engine.", 18, MUTED, False)]])
notes(s, """
The new way. She says it in plain language, the way she'd ask a person. The agent
already has a calendar skill — a small text file of instructions for how it uses
its calendar tools. It appends one line. That's the entire implementation. No
recipient picker, no trigger UI, no backend rule. The behavior she wanted now
exists, and it was authored by the agent, on the spot, in the channel she was
already using.
""", tech="""
TWO LAYERS — the authoring, and the firing.

1) THE AUTHORING (fully working):
• The skill is skills/calendar/SKILL.md — YAML frontmatter (name + description,
  ≤1024 chars, injected into the system prompt as Level-1 metadata) plus a ≤5 KB
  markdown body of instructions.
• The agent edits it at runtime from inside the sandbox via the bb.skill module
  (create / set body / save → file written boxbot:boxbot, 0644). load_skill pulls
  the body into context on demand (3-level progressive disclosure).
• So "one line added" is real: a diff to one markdown file, authored by the agent.

2) THE FIRING — BE HONEST IF ASKED (designed; today via polling, not a webhook):
• Working today: triggers (core/scheduler.py — fire_at / fire_after / cron /
  person, AND-combinable); the agent wakes on a TriggerFired event
  (core/agent.py::_on_trigger_fired) and runs a one-shot conversation; outbound
  WhatsApp via the `message` tool → core/output_dispatcher.py → MessageRouter →
  Meta Graph API (communication/{router.py, whatsapp.py}).
• NOT yet wired: a Google Calendar push-webhook, or a "poll-and-diff" trigger type.
• Honest demo framing: a cron trigger (e.g. every ~10 min) calls the calendar
  integration, diffs against a snapshot in bb.workspace, and texts on a new event.
  You can present this as "works today via polling" — just don't imply an instant
  event-driven push, which isn't built yet.
""")

# =============================================================================
# 7 — WHAT JUST HAPPENED
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "What just happened")
pairs = [
    ("The interface", "A rule-builder UI", "One sentence of plain English"),
    ("The implementation", "Channel + recipients + trigger + conditions", "One line in a text file"),
    ("Who built it", "Our engineering org, in advance", "The agent, on demand"),
]
y = Inches(2.0)
colx_label = MX
colx_old = Inches(4.3)
colx_new = Inches(8.7)
text(s, colx_old, y - Inches(0.55), Inches(4), Inches(0.4), [[("OLD MODEL", 12, MUTED, True, 2)]])
text(s, colx_new, y - Inches(0.55), Inches(4), Inches(0.4), [[("AGENT MODEL", 12, AMBER, True, 2)]])
for label, old, new in pairs:
    panel(s, MX - Inches(0.15), y, Inches(12.4), Inches(1.3), fill=PANEL)
    text(s, colx_label, y, Inches(3.4), Inches(1.3), [[(label, 17, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, colx_old, y, Inches(4.2), Inches(1.3), [[(old, 15, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, colx_new, y, Inches(4.2), Inches(1.3), [[(new, 15, AMBER, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.55)
text(s, MX, Inches(6.8), Inches(11.5), Inches(0.5),
     [[("Application logic gets compressed. The desired behavior is reached with almost no effort.", 16, INK, True)]])
notes(s, """
Same job, three columns of difference. The interface collapsed from a builder to a
sentence. The implementation collapsed from a four-part rule to one line. And the
work moved from us-in-advance to the-agent-on-demand. For a CX org this is the
headline: the gap between 'what the customer wants' and 'the product does it'
shrinks to one message. That gap is where churn and support tickets live today.
""", tech="""
• Concretely, the "implementation" is a diff to one markdown file. No schema
  migration, no deploy, no feature flag, no app-store review.
• It was authored by the agent from inside the OS-level sandbox (boxbot-sandbox
  user; seccomp blocks execve/fork; .env is 0600 and unreadable to the sandbox;
  site-packages read-only). So letting users drive changes via the agent does NOT
  hand them a path into core code or secrets — the blast radius is a text file.
• The delivery path it "wrote" already existed (the `message` tool). The agent
  composed existing primitives; it didn't generate privileged code. That's the
  safety story behind "the agent builds features."
""")

# =============================================================================
# 8 — MEET boxBot (the device / feature set) + demo
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "Meet boxBot")
text(s, MX, Inches(1.6), Inches(6.4), Inches(1.4),
     [[("An executive assistant", 30, INK, True)],
      [("for the household.", 30, INK, True)]], space_after=2)
feats = [
    ("Sees", "recognizes the family through the camera"),
    ("Hears & speaks", "walk up and talk — mic + speaker"),
    ("Remembers", "people, preferences, and context over time"),
    ("Manages", "calendar, tasks, and the photo library"),
    ("Reaches you", "texts over WhatsApp, anywhere"),
]
y = Inches(3.25)
for name, desc in feats:
    text(s, MX, y, Inches(2.5), Inches(0.5), [[(name, 16, AMBER, True)]])
    text(s, MX + Inches(2.55), y, Inches(3.7), Inches(0.5), [[(desc, 14, MUTED, False)]])
    y += Inches(0.62)
placeholder(s, Inches(7.5), Inches(1.6), Inches(4.9), Inches(5.2),
            "DEMO CLIP  ·  ~45s", "live: “BB, what's on today?” → voice reply + screen updates")
notes(s, """
Now the device itself. Raspberry Pi in a wooden box: camera, mic, speaker, a 7-inch
screen, and a WhatsApp line. Functionally it's a family executive assistant — it
knows who's in the room, holds a conversation, remembers what matters, runs the
calendar and tasks and photos, and can reach you when you're out. [PLAY THE DEMO
CLIP HERE] — keep it short: one spoken request, the voice answer, and the screen
reacting. Let them feel that it's a real, present thing, not a chatbot.
""", tech="""
VOICE PATH (src/boxbot/communication/): OpenWakeWord "BB" (ONNX) → Silero VAD →
AudioCapture (utterance boundaries) → pyannote diarization on CPU (attributed,
multi-speaker) → ElevenLabs Scribe STT → agent → ElevenLabs streaming TTS.
~3-4s from end-of-speech to first spoken word. NOTE: not a real-time streaming
voice agent — the agent receives attributed transcripts and decides whether/when
to reply.

VISION / WHO'S IN THE ROOM (src/boxbot/perception/): YOLOv5s-personface on the
Hailo NPU → RepVGG-A0 ReID embeddings (128-dim) on Hailo; pyannote voice
embeddings + mic-array DOA fused in fusion.py. Voice-confirmed identity gates
whether visual embeddings get written to the cloud. Identities are injected into
the prompt as tiered text, e.g. "[Jacob] (voice: high 0.92)" — the agent never
sees raw embeddings.

ONE CONVERSATION ABSTRACTION (core/conversation.py): states LISTENING / THINKING /
SPEAKING / ENDED. Voice = transient (180s silence timer, in-memory). WhatsApp =
persistent (SQLite ConversationStore, ~4h window). Voice and WhatsApp run as
separate conversations in parallel.

HONEST FLAG (barge-in): the design mentions "graduated yielding," but what's built
is binary, wake-word-gated re-engagement — STT detaches from the mic during TTS;
saying "BB" cancels the in-flight generation, folds the partial reply in as an
interrupted turn, and restarts. Cleaner for noisy multi-speaker rooms; just don't
claim soft/volume-fade barge-in.
""")

# =============================================================================
# 9 — HOW IT WORKS: TOOLS / SKILLS / DISPLAYS
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "The design philosophy")
text(s, MX, Inches(1.6), Inches(11.5), Inches(0.8),
     [[("The agent has a small set of ", 28, INK, False),
       ("hands", 28, AMBER, True),
       (", and text files that tell it how to use them.", 28, INK, False)]])
cards = [
    ("TOOLS", "The hands", "Always attached. Core actions the agent can always take — switch the screen, send a message, recall a memory."),
    ("SKILLS", "What it picks up", "Modular text-file capabilities. Weather, reminders, the calendar. A new feature is a new file."),
    ("DISPLAYS", "What it shows", "The screen is the agent's canvas, not a fixed widget set. It builds the layout the moment calls for."),
]
cw = Inches(3.78); gap = Inches(0.28); x = MX; y = Inches(2.95); ch = Inches(3.5)
for title, sub, body in cards:
    panel(s, x, y, cw, ch, fill=PANEL)
    bar(s, x + Inches(0.35), y + Inches(0.4), Inches(0.45), Pt(3), AMBER)
    text(s, x + Inches(0.35), y + Inches(0.6), cw - Inches(0.6), Inches(0.5),
         [[(title, 16, AMBER, True, 2)]])
    text(s, x + Inches(0.35), y + Inches(1.1), cw - Inches(0.6), Inches(0.5),
         [[(sub, 21, INK, True)]])
    text(s, x + Inches(0.35), y + Inches(1.75), cw - Inches(0.65), Inches(1.6),
         [[(body, 14.5, MUTED, False)]], line_spacing=1.12)
    x += cw + gap
notes(s, """
Under the hood it's deliberately simple, and this is the part to internalize. The
agent has a SMALL number of tools — its hands, always attached. It has SKILLS —
modular text files that teach it how to use those hands for a given job; that
calendar line from earlier lived in a skill. And it has DISPLAYS — the screen is
something it composes, not a fixed set of widgets we shipped. Three primitives.
Almost everything the product 'does' is now text the agent reads and edits, not
code we hard-wire.
""", tech="""
• Slight correction for accuracy if asked: there are actually 10 always-loaded
  tools (tools/registry.py): execute_script, switch_display, identify_person,
  manage_tasks, search_memory, search_photos, web_search, load_skill, message,
  mute_mic. The slide says "small set" — that's the honest spirit.
• execute_script is the gateway. It runs Python as the boxbot-sandbox user and
  speaks streaming bidirectional JSON: the sandbox emits __BOXBOT_SDK_ACTION__
  lines, the main process dispatches each to a per-module handler
  (tools/_sandbox_actions.py), writes a JSON reply back, and collects any image
  attachments into a multimodal tool result. That's how one script can compose
  many ops (camera + workspace + display + memory…) in a single turn.
• The agent's own text output is a PRIVATE scratchpad (INTERNAL_NOTES_SCHEMA:
  thought / observations — never delivered). To reach a human it MUST call the
  `message` tool. That structural boundary is why nothing leaks by accident.
• Skills = skills/*/SKILL.md (3-level progressive disclosure). Displays =
  declarative JSON specs (block tree only, no executable render code) — so there's
  no privileged code path to gate when the agent authors one.
""")

# =============================================================================
# 10 — DISPLAYS: the agent builds the screen (real screenshots)
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.8), "Displays, not widgets")
text(s, MX, Inches(1.45), Inches(11.5), Inches(0.7),
     [[("The agent composes the screen for the moment.", 27, INK, True)]])
shots = [
    ("morning_brief_boxbot.png", "Morning brief"),
    ("home_metrics_boxbot.png", "Home & power"),
    ("now_listening_boxbot.png", "In conversation"),
]
cw = Inches(3.85); gap = Inches(0.18); x = MX; y = Inches(2.55)
imh = Inches(2.55)
for fn, cap in shots:
    image_contain(s, os.path.join(IMG, fn), x, y, cw, imh)
    text(s, x, y + imh + Inches(0.12), cw, Inches(0.4),
         [[(cap, 14, MUTED, False)]], align=PP_ALIGN.CENTER)
    x += cw + gap
text(s, MX, Inches(5.75), Inches(11.5), Inches(1.0),
     [[("Same screen, three layouts — chosen by context, not by a settings toggle. ", 16, MUTED, False),
       ("These are live boxBot screens.", 16, AMBER, True)]])
notes(s, """
These are actual screens off the device. The morning brief, a home-and-power
dashboard, the in-conversation view. The user never picked these from a gallery and
never arranged the tiles. The agent assembled each one for the moment it was needed
— morning vs. mid-conversation. The screen is fluid. Which sets up the obvious
question: what if the user wants something we never designed?
""", tech="""
• Block system (src/boxbot/displays/): 24 block types — 7 layout containers (row,
  column, columns, card, spacer, divider, repeat), 13 content blocks (text, metric,
  badge, list, table, key_value, icon, emoji, image, chart, progress, clock,
  countdown), 2 composite widgets (weather, calendar), 2 meta (rotate, page_dots).
  Registered in blocks.py::BLOCK_REGISTRY.
• A display is a JSON spec (see displays/morning_brief/display.json): name, theme,
  data_sources[], and one layout tree. No code in the spec.
• Live data via {source.field} bindings, resolved in spec.py. A DataSourceManager
  (data_sources.py) runs a per-source async refresh loop and caches last-good data,
  so a fetch error shows stale values, never a blank tile. Renderer (renderer.py)
  composes a PIL 1024x600 image → pygame.
• These three are real preview renders off that pipeline. The agent authors a spec
  as a dict and calls bb.display.preview → it VIEWS the resulting PNG (multimodal)
  → iterates → bb.display.save. It never writes render code.
""")

# =============================================================================
# 11 — DISPLAYS ON DEMAND
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "“Add that to my dashboard.”")
text(s, MX, Inches(1.75), Inches(11.5), Inches(1.4),
     [[("Want something we never shipped?", 30, INK, True)],
      [("The agent builds the data pipe and plugs it into a widget.", 22, MUTED, False)]],
     space_after=6)
chips = ["Solar generation", "Surf conditions", "Sports headlines",
         "Flight status", "Pollen count", "Stock ticker"]
per_row = 3
ch_w = Inches(3.6); ch_h = Inches(0.78); gap = Inches(0.35)
row_w = per_row * ch_w + (per_row - 1) * gap
x0 = MX + (Inches(11.5) - row_w) // 2
y0 = Inches(3.9)
for i, c in enumerate(chips):
    r = i // per_row; col = i % per_row
    x = x0 + col * (ch_w + gap)
    y = y0 + r * (ch_h + Inches(0.3))
    panel(s, x, y, ch_w, ch_h, fill=PANEL, line=AMBER_DK)
    text(s, x, y, ch_w, ch_h, [[(c, 16, INK, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, MX, Inches(6.5), Inches(11.5), Inches(0.6),
     [[("No backlog ticket. No release. The capability the user asked for simply appears.", 16, AMBER, True)]])
notes(s, """
This is the unlock. 'Add solar generation to my summary.' 'Put surf conditions on
the morning screen.' In today's world each of those is a feature request that joins
a backlog. Here, the agent finds or builds the data pipeline, maps it onto an
existing widget, and the tile appears. The widget LIBRARY is the durable thing we
build once; the infinite arrangements of it come for free. That's the leverage.
""", tech="""
• The "data pipe" is concrete: either an INTEGRATION the agent authors
  (bb.integrations.create → manifest.yaml + script.py, sandbox-run on call), or a
  built-in http_json data source with declarative field extraction and value→
  icon/color `map` transforms — no code in the display spec either way.
• Any API key goes through bb.secrets, which is WRITE-ONLY to the sandbox: the
  agent can store a key but can never read it back. The main process injects it as
  a BOXBOT_SECRET_<NAME> env var only when the declaring integration runs. So the
  agent can USE a credential it cannot exfiltrate.
• A new integration/source is live on the next read — no restart. The durable
  asset we build once is the widget/block + source LIBRARY; every arrangement of it
  (solar, surf, sports…) is then free and agent-assembled.
• Honest aside: bb.packages (request a new pip dependency, human-approved) is
  designed but the approval handler isn't wired yet — so "new data source" today
  means stdlib/requests-style fetches, not arbitrary new packages.
""")

# =============================================================================
# 12 — AGENT-FIRST ARCHITECTURE (dark mode story)
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "Agent-first architecture")
text(s, MX, Inches(1.7), Inches(11.5), Inches(1.6),
     [[("Settings, layouts, and behaviors aren't hard-coded.", 30, INK, True)],
      [("They're configuration the agent reads and edits.", 24, MUTED, False)]],
     space_after=8)
# dark mode mini-flow
steps = [
    ("Today", "“Add dark mode” = a multi-quarter project, thousands of tickets", MUTED),
    ("Agent model", "User: “I prefer dark mode.”  Agent edits the theme.  “Does this look right?”", AMBER),
]
y = Inches(4.0)
for label, body, col in steps:
    panel(s, MX, y, Inches(11.5), Inches(1.15), fill=PANEL)
    text(s, MX + Inches(0.4), y, Inches(2.4), Inches(1.15), [[(label, 17, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, MX + Inches(3.0), y, Inches(8.2), Inches(1.15), [[(body, 16, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.35)
notes(s, """
Step up a level from displays to the whole app. If layout, theme, and behavior are
configuration the agent can read and write — not code compiled into the binary —
then 'changes to the product' stop being engineering projects. Dark mode is the
canonical example. Today it's a multi-quarter effort with a thousand tickets. In
this architecture the user says 'I prefer dark mode,' the agent edits the theme
value, and asks 'does this look right?' Up-front work to build the harness; an
enormous amount of capability for free afterward.
""", tech="""
• Themes are config, not code: 4 built-ins as dataclasses in displays/themes.py
  (boxbot, midnight, daylight, classic), plus community themes as YAML in themes/.
  Each defines colors / fonts / spacing / radius / shadow / icon_style.
• "Add dark mode" = the agent setting a spec's `theme` field (or editing theme
  values) → bb.display.preview → it views the PNG → asks the user. Same pattern for
  layout (edit the block tree) and behavior (edit a SKILL.md or system.md line).
• Guardrails make this safe to expose: specs are validated before save (spec.py),
  and system memory has hard limits (4 KB cap, fixed section schema, version
  history, no-secrets rule). The agent edits configuration; it can't reach in and
  rewrite the application.
""")

# =============================================================================
# 13 — WHY THIS MATTERS (for CX / Alarm.com)
# =============================================================================
s = slide()
kicker(s, MX, Inches(0.95), "Why this matters")
cards = [
    ("Experience gets stickier", "Walk up and talk. Once it runs the home, the calendar, and the photos, the customer never leaves — they take it with them when they move."),
    ("Engineering burden collapses", "No menus, rules engines, or onboarding wizards to maintain. A new feature becomes a new text file. The coding footprint scales down."),
    ("Trust becomes the moat", "A device the family texts every day builds a relationship a bolted-on chatbot never will. Familiarity and reliability are the lock-in."),
]
y = Inches(1.95)
for title, body in cards:
    panel(s, MX, y, Inches(11.5), Inches(1.45), fill=PANEL)
    bar(s, MX, y + Inches(0.28), Pt(4), Inches(0.9), AMBER)
    text(s, MX + Inches(0.45), y + Inches(0.22), Inches(10.6), Inches(0.5), [[(title, 20, INK, True)]])
    text(s, MX + Inches(0.45), y + Inches(0.72), Inches(10.6), Inches(0.6), [[(body, 14.5, MUTED, False)]])
    y += Inches(1.65)
notes(s, """
Bring it home for a CX audience. Three layered wins. One, the experience is better
and stickier — walk-up-and-talk, and once it's holding the household's life it
moves with the customer. Two, our build-and-maintain cost drops hard: no more
sprawling settings, rule engines, and wizards to keep alive — features become text
files. Three, and most strategic: trust. A thing your family talks to daily earns a
relationship a website chatbot can't. In the AI age, that relationship plus the
backend is the moat.
""", tech="""
• Stickiness is backed by memory (src/boxbot/memory/, data/memory/): an
  always-loaded system.md of household facts + typed fact memories in SQLite,
  retrieved by hybrid search (on-device MiniLM 384-dim vectors, 0.6, + BM25, 0.4)
  with a Haiku rerank. It accrues a model of the household and carries between
  rooms/devices — that's the "knows you" feeling.
• Cost: the two-model split keeps the cheap model (Haiku) on the high-frequency,
  low-stakes work. Your thesis figure is ~$20/mo/home and falling; the task needs
  bounded intelligence, and the Hailo NPU absorbs simple on-device work.
• "Engineering burden collapses" is the OS-sandbox story in reverse: because
  agent-authored changes are just text files in an isolated sandbox (seccomp blocks
  execve/fork, .env 0600 unreadable, site-packages read-only), you can let the
  product expand via the agent without expanding the attack surface or the
  maintenance tree.
""")

# =============================================================================
# 14 — CLOSE
# =============================================================================
s = slide()
bar(s, MX, Inches(2.5), Inches(0.7), Pt(4), AMBER)
text(s, MX, Inches(2.8), Inches(11.5), Inches(2.4),
     [[("The UI becomes", 50, INK, True)],
      [("whatever the user wants.", 50, AMBER, True)]],
     space_after=6, line_spacing=1.04)
text(s, MX, Inches(5.2), Inches(11.5), Inches(0.7),
     [[("Stop shipping surfaces. Ship an assistant that builds them.", 22, MUTED, False)]])
text(s, MX, Inches(6.6), Inches(11.5), Inches(0.5),
     [[("boxBot   ·   Smart Home 2.0", 14, FAINT, True, 2)]])
notes(s, """
Close on the inversion we opened with. We've spent a decade getting very good at
building surfaces — screens, settings, rules. The opportunity is to stop shipping
surfaces and ship the assistant that builds them on demand. We're one of the few
companies actually positioned to do it: the dealer channel, the device integrations,
the security backend the FAANG players can't replicate. I'd love to talk about
where CX leads from here. Questions.
""", tech="""
• The summary architecture point: the durable, hand-built assets are SMALL and
  stable — the agent harness (loop + Conversation abstraction), the 10 tools, the
  `bb` sandbox SDK, the widget/block + theme library, and the backend integrations.
  Everything user-facing (skills, displays, behaviors) is generated text/specs.
• That's the whole thesis in one line: a feature is a text file, and the UI is
  whatever the user asks for, because the agent composes durable primitives instead
  of us shipping new surfaces.
• Honest "not-yet" list for credibility if pushed: calendar firing is poll-and-diff
  not a push-webhook; barge-in is binary wake-word, not graduated; bb.packages
  approval flow isn't wired. All are designed-for and small to close — good
  "here's the roadmap" answers, not gaps in the core idea.
""")

out = os.path.join(HERE, "boxBot_SmartHome_2.0.pptx")
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
